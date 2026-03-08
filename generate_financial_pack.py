import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import  XL_LABEL_POSITION
from pptx.util import Pt
import datetime
from pptx.dml.color import RGBColor
import numpy as np


# -----------------------------
# Helper function to clean dataframe columns
# -----------------------------
def clean_columns(df):
    df.columns = df.columns.str.strip().str.lower()
    return df

# Function to clean NaN and INF values
def clean_df(df):
    df = df.replace([np.inf, -np.inf], 0)
    df = df.fillna(0)
    return df


# -----------------------------
# Load CSV files safely
# -----------------------------
pnl = clean_columns(pd.read_csv("vw_pnl_final.csv"))
waterfall = clean_columns(pd.read_csv("vw_pnl_waterfall.csv"))
yoy = clean_columns(pd.read_csv("vw_yoy_variance.csv"))
budget = clean_columns(pd.read_csv("vw_budget_vs_actual.csv"))
expense = clean_columns(pd.read_csv("vw_expense_positive.csv"))
ytd = clean_columns(pd.read_csv("vw_ytd_summary.csv"))
df = pd.read_csv("gl_transaction_with_supplier.csv")
variance_q1 = pd.read_csv("vw_final_variance_q1_2025.csv")
variance_year = pd.read_csv("vw_full_year_2025_variance.csv")
waterfall = pd.read_csv("vw_variance_waterfall_2025_q1.csv")
accrual = pd.read_csv("vw_accrual_report.csv")
accrual = pd.read_csv("vw_accrual_report.csv")      # accrual report
circuit = pd.read_csv("circuit_list.csv")
fx = pd.read_csv("vw_fx_revaluation_dashboard_clean.csv")
debtor_aging = pd.read_csv("vw_net_debtor_aging_summary.csv")
cap_df = pd.read_csv("vw_asset_capitalization.csv")
dep_df = pd.read_csv("vw_asset_depreciation.csv")
actual_2025 = pd.read_csv("vw_actuals_2025_with_journals.csv")
actual_2026 = pd.read_csv("vw_actuals_2026.csv")
growth_df = pd.read_csv("vw_growth_trend_realistic.csv")
forecast_df = pd.read_csv("vw_budget_forecast_2027.csv")
exp_variance = clean_columns(pd.read_csv("vw_expense_variance.csv"))
exp_trend = clean_columns(pd.read_csv("vw_expense_trend.csv"))
exp_comparison = clean_columns(pd.read_csv("vw_expense_comparison.csv"))
depr_forecast = clean_columns(pd.read_csv("monthly_total_depr.csv"))
cwip_summary = clean_columns(pd.read_csv("vw_cwip_summary.csv"))


# Clean column names
fx.columns = fx.columns.str.strip().str.lower()

# Convert month_end to date
fx['month_end'] = pd.to_datetime(fx['month_end'])


# Fill missing values
pnl = pnl.fillna(0)
yoy = yoy.fillna(0)
budget = budget.fillna(0)
expense = expense.fillna(0)
ytd = ytd.fillna(0)
actual_2025 = actual_2025.replace([np.inf, -np.inf], 0).fillna(0)
actual_2026 = actual_2026.replace([np.inf, -np.inf], 0).fillna(0)
growth_df = growth_df.replace([np.inf, -np.inf], 0).fillna(0)
forecast_df = forecast_df.replace([np.inf, -np.inf], 0).fillna(0)
cap_df = clean_df(cap_df)
dep_df = clean_df(dep_df)
exp_variance = clean_df(exp_variance)
exp_trend = clean_df(exp_trend)
exp_comparison = clean_df(exp_comparison)
exp_comparison.fillna(0, inplace=True)
depr_forecast.fillna(0, inplace=True)
cwip_summary.fillna(0, inplace=True)


# Pivot Debtor Aging
# -----------------------------
pivot_aging = debtor_aging.pivot(
    index='operator_name',
    columns='aging_bucket',
    values='total_net_balance'
).fillna(0)

bucket_order = ['0-30 Days', '31-60 Days', '61-90 Days', '90+ Days']
pivot_aging = pivot_aging.reindex(columns=bucket_order, fill_value=0)

# -----------------------------
# Create Presentation
# -----------------------------
prs = Presentation()
blank_layout = prs.slide_layouts[5]

# -----------------------------
# 1. Title Slide
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Auto Financial Pack - FY 2025"

# -----------------------------
# 2. Revenue & Profit Trend
# -----------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Revenue & Profit Trend"

chart_data = CategoryChartData()
chart_data.categories = pnl['period'].astype(str).tolist()
chart_data.add_series('Revenue', pnl['revenue'].tolist())
chart_data.add_series('EBITDA', pnl['ebitda'].tolist())
chart_data.add_series('Profit After Tax', pnl['profit_after_tax'].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------
# 3. YOY Growth
# -----------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Year Over Year Growth %"

chart_data = CategoryChartData()
chart_data.categories = yoy['quarter'].astype(str).tolist()
chart_data.add_series('YOY Growth %', yoy['yoy_growth_percentage'].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------
# 4. Budget vs Actual
# -----------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Budget vs Actual"

chart_data = CategoryChartData()
chart_data.categories = budget['pnl_line'].astype(str).tolist()
chart_data.add_series('Actual', budget['actual_amount'].tolist())
chart_data.add_series('Budget', budget['budget_amount'].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------
# 5. Expense Breakdown
# -----------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Expense Breakdown"

chart_data = CategoryChartData()
chart_data.categories = expense['period'].astype(str).tolist()
chart_data.add_series('Opex', expense['opex'].tolist())
chart_data.add_series('Depreciation', expense['depreciation'].tolist())
chart_data.add_series('Tax', expense['tax'].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------
# 6. YTD Summary
# -----------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "YTD Summary"

textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(3))
tf = textbox.text_frame

tf.text = f"YTD Revenue: {ytd['ytd_revenue'].iloc[0]:,.2f}"

p = tf.add_paragraph()
p.text = f"YTD EBITDA: {ytd['ytd_ebitda'].iloc[0]:,.2f}"

p = tf.add_paragraph()
p.text = f"YTD Operating Profit: {ytd['ytd_operating_profit'].iloc[0]:,.2f}"

p = tf.add_paragraph()
p.text = f"YTD Profit After Tax: {ytd['ytd_profit_after_tax'].iloc[0]:,.2f}"


# =====================================================
# 7. Product Revenue Comparison (2025)
# =====================================================
product_rev = clean_columns(pd.read_csv("vw_product_revenue_comparison.csv"))
product_rev =_finalize = product_rev.fillna(0)

# Filter latest year (2025)
product_rev_2025 = product_rev[product_rev['year'] == 2025]

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Product Revenue Comparison - 2025"

chart_data = CategoryChartData()
chart_data.categories = product_rev_2025['product_line']

chart_data.add_series(
    'Revenue',
    product_rev_2025['revenue'].tolist()
)

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# =====================================================
# 8. Product Profit Comparison (2025)
# =====================================================
product_profit = clean_columns(pd.read_csv("vw_product_profit_comparison.csv"))
product_profit = product_profit.fillna(0)

product_profit_2025 = product_profit[product_profit['year'] == 2025]

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Product Profit Comparison - 2025"

chart_data = CategoryChartData()
chart_data.categories = product_profit_2025['product_line']

chart_data.add_series(
    'Profit',
    product_profit_2025['profit'].tolist()
)

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# =====================================================
# 9. Product Profit Margin %
# =====================================================
product_margin = clean_columns(pd.read_csv("vw_product_margin_comparison.csv"))
product_margin = product_margin.fillna(0)

product_margin_2025 = product_margin[product_margin['year'] == 2025]

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Product Profit Margin % - 2025"

chart_data = CategoryChartData()
chart_data.categories = product_margin_2025['product_line']

chart_data.add_series(
    'Profit Margin %',
    (product_margin_2025['profit_margin'] * 100).tolist()
)

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# =====================================================
# 10. Product YOY Growth %
# =====================================================
product_yoy = clean_columns(pd.read_csv("vw_product_yoy.csv"))
product_yoy = product_yoy.fillna(0)

product_yoy_2025 = product_yoy[product_yoy['year'] == 2025]

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Product YOY Revenue Growth % - 2025"

chart_data = CategoryChartData()
chart_data.categories = product_yoy_2025['product_line']

chart_data.add_series(
    'Product YOY Growth %',
    (product_yoy_2025['yoy_growth_percent'] * 100).tolist()
)

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# 11. Top 5 Suppliers (5-Year Total)
# -----------------------------
top5 = (
    df.groupby("supplier_name")["amount"]
    .sum()
    .reset_index()
    .sort_values(by="amount", ascending=False)
    .head(5)
)

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
).text = "Top 5 Suppliers - 5 Year Advertising Spend"

chart_data = CategoryChartData()
chart_data.categories = top5["supplier_name"]
chart_data.add_series("Total Spend", top5["amount"].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------
# 12. Total Advertising Spend
# -----------------------------
total_spend = df["amount"].sum()

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
).text = "Total Advertising Expenditure"

textbox = slide.shapes.add_textbox(
    Inches(2), Inches(2), Inches(6), Inches(3)
)

textbox.text_frame.text = f"Total 5-Year Advertising Spend: {total_spend:,.2f}"

# -----------------------------
# 13. Monthly Trend
# -----------------------------
df["transaction_date"] = pd.to_datetime(df["transaction_date"])
df["year_month"] = df["transaction_date"].dt.to_period("M").astype(str)

monthly = (
    df.groupby("year_month")["amount"]
    .sum()
    .reset_index()
)

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
).text = "Monthly Advertising Trend"

chart_data = CategoryChartData()
chart_data.categories = monthly["year_month"]
chart_data.add_series("Monthly Spend", monthly["amount"].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------
# 14. Quarterly Trend
# -----------------------------
quarterly = (
    df.groupby(["year", "quarter"])["amount"]
    .sum()
    .reset_index()
)

quarterly["year_quarter"] = quarterly["year"].astype(str) + "-" + quarterly["quarter"]

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
).text = "Quarterly Advertising Trend"

chart_data = CategoryChartData()
chart_data.categories = quarterly["year_quarter"]
chart_data.add_series("Quarterly Spend", quarterly["amount"].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------
# 15. Supplier Contribution %
# -----------------------------
contribution = (
    df.groupby("supplier_name")["amount"]
    .sum()
    .reset_index()
)

total = contribution["amount"].sum()
contribution["percentage"] = (contribution["amount"] / total) * 100

slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
).text = "Supplier Contribution %"

chart_data = CategoryChartData()
chart_data.categories = contribution["supplier_name"]
chart_data.add_series("Contribution %", contribution["percentage"].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(2), Inches(1),
    Inches(6), Inches(5),
    chart_data
)

# -----------------------------------
# 16. Full Year 2025 Variance
# -----------------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.title.text = "Full Year 2025 Variance"

chart_data = CategoryChartData()
chart_data.categories = variance_year["pnl_line"]
chart_data.add_series("Variance Amount", variance_year["variance_amount"])

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4.5),
    chart_data
)

# -----------------------------------
# 17. Q1 2025 Waterfall
# -----------------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.title.text = "Q1 2025 Variance Waterfall"

chart_data = CategoryChartData()
chart_data.categories = waterfall["metric"]
chart_data.add_series("Variance", waterfall["value"])

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,  # python-pptx doesn't support real waterfall
    Inches(1), Inches(1.5),
    Inches(8), Inches(4.5),
    chart_data
)

# ---------------------------
# 18. Accrual Report
# ---------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Accrual Report"

# 19. Standardize column names
accrual.columns = accrual.columns.str.strip().str.lower()
df.columns = df.columns.str.strip().str.lower()
circuit.columns = circuit.columns.str.strip().str.lower()

# 20. Total accrual and expense
total_accrual = accrual['accrual_amount'].sum()
total_expense = df['amount'].sum()
accrual_percent = (total_accrual / total_expense) * 100 if total_expense != 0 else 0

# 21. Add textbox summary
textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2.5))
tf = textbox.text_frame
tf.text = f"Total Accrual Amount: {total_accrual:,.0f}"
p = tf.add_paragraph()
p.text = f"Total Expense Amount: {total_expense:,.0f}"
p = tf.add_paragraph()
p.text = f"Accrual % of Total Expense: {accrual_percent:.2f}%"

# 22. Top 5 suppliers by accrual
top_suppliers = accrual.groupby('supplier_name')['accrual_amount'].sum().sort_values(ascending=False).head(5)
p = tf.add_paragraph()
p.text = "Top 5 Suppliers by Accrual:"
for supplier, amt in top_suppliers.items():
    p = tf.add_paragraph()
    p.text = f"{supplier}: {amt:,.0f}"

# 23. Accrual per GL code summary (needed for chart)
accrual_summary = accrual.groupby('gl_code')['accrual_amount'].sum().reset_index()
accrual_summary = accrual_summary.rename(columns={'accrual_amount':'total_accrual_amount'})

# Circuit summary per supplier (needed for chart)
circuit_summary = circuit.groupby('supplier_name')['amount'].sum().reset_index()
circuit_summary = circuit_summary.rename(columns={'amount':'total_circuit_amount'})

# ---------------------------
# 24. Accrual / Circuit Slide with Charts
# ---------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Accrual / Circuit Report"

# 25. Textbox with totals
textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
tf = textbox.text_frame
tf.text = f"Total Accrual: {total_accrual:,.0f}"
p = tf.add_paragraph()
p.text = f"Total Circuit Amount: {circuit['amount'].sum():,.0f}"
p = tf.add_paragraph()
p.text = f"Accrual %: {accrual_percent:.2f}%"

# Chart 1: Accrual per GL code
chart_data = CategoryChartData()
chart_data.categories = accrual_summary['gl_code'].tolist()
chart_data.add_series('Accrual Amount', accrual_summary['total_accrual_amount'].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(3),
    Inches(4), Inches(3.5),
    chart_data
)

# Chart 2: Circuit per Supplier
chart_data2 = CategoryChartData()
chart_data2.categories = circuit_summary['supplier_name'].tolist()
chart_data2.add_series('Circuit Amount', circuit_summary['total_circuit_amount'].tolist())

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(5), Inches(3),
    Inches(4), Inches(3.5),
    chart_data2
)

# ---------------------------
# 26. FX REVALUATION REPORT
# ---------------------------

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)
).text = "FX Revaluation Report"

# 27. Summary numbers
total_balance_lkr = fx['balance_lkr'].sum()
total_gain_loss = fx['reval_gain_loss'].sum()

textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1), Inches(8), Inches(2))
tf = textbox.text_frame
tf.text = f"Total Balance (LKR): {total_balance_lkr:,.0f}"

p = tf.add_paragraph()
p.text = f"Total FX Revaluation Gain/Loss: {total_gain_loss:,.0f}"

# 28. Top suppliers by FX impact
top_fx = fx.groupby('supplier_name')['reval_gain_loss'].sum().sort_values(ascending=False).head(5)

p = tf.add_paragraph()
p.text = "Top 5 Suppliers by FX Gain/Loss:"

for supplier, value in top_fx.items():
    p = tf.add_paragraph()
    p.text = f"{supplier}: {value:,.0f}"

# Chart 1: FX Gain/Loss by Supplier
fx_supplier = fx.groupby('supplier_name')['reval_gain_loss'].sum().reset_index()

chart_data = CategoryChartData()
chart_data.categories = fx_supplier['supplier_name'].tolist()
chart_data.add_series(
    'FX Gain/Loss',
    fx_supplier['reval_gain_loss'].tolist()
)

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8), Inches(3.5),
    Inches(4.2), Inches(3),
    chart_data
)


# Chart 2: Monthly FX Movement
fx_month = fx.groupby('month_end')['reval_gain_loss'].sum().reset_index()

chart_data2 = CategoryChartData()
chart_data2.categories = fx_month['month_end'].dt.strftime('%Y-%m').tolist()
chart_data2.add_series(
    'Monthly FX Gain/Loss',
    fx_month['reval_gain_loss'].tolist()
)

slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(5.2), Inches(3.5),
    Inches(4.2), Inches(3),
    chart_data2
)

# 29. Debtor Aging Slide
# -----------------------------
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Net Debtor Aging Report"

# 30. Summary Textbox
textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
tf = textbox.text_frame
tf.text = "Debtor Aging Summary per Operator"
for operator in pivot_aging.index:
    p = tf.add_paragraph()
    p.text = f"{operator}: " + ", ".join([f"{bucket}: {pivot_aging.loc[operator, bucket]:,.0f}" for bucket in bucket_order])

# 31. Stacked Column Chart
chart_data = CategoryChartData()
chart_data.categories = pivot_aging.index.tolist()
for bucket in pivot_aging.columns:
    chart_data.add_series(bucket, pivot_aging[bucket].tolist())

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED,
    Inches(1), Inches(3),
    Inches(8), Inches(4.5),
    chart_data
).chart

# 32. Apply consistent colors per bucket
bucket_colors = {
    '0-30 Days': RGBColor(91, 155, 213),    # Blue
    '31-60 Days': RGBColor(237, 125, 49),   # Orange
    '61-90 Days': RGBColor(165, 165, 165),  # Gray
    '90+ Days': RGBColor(255, 192, 0)       # Yellow
}

for i, series in enumerate(chart.series):
    bucket_name = pivot_aging.columns[i]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = bucket_colors[bucket_name]

# -----------------------------
# 33 .Define colors per asset group
# -----------------------------
asset_colors = {
    'Laptop': RGBColor(91, 155, 213),   # Blue
    'Router': RGBColor(237, 125, 49),   # Orange
    'Printer': RGBColor(165, 165, 165)  # Gray
}

# -----------------------------
# 34. Asset Capitalization by Asset Group
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Asset Capitalization by Asset Group"

chart_data = CategoryChartData()
chart_data.categories = cap_df['asset_group']
chart_data.add_series('Total Capitalized Value', cap_df['total_capitalized_value'])

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4),
    chart_data
).chart

# 35. Apply colors correctly
for series in chart.series:
    for j, point in enumerate(series.points):
        category_name = chart_data.categories[j].label
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = asset_colors.get(category_name, RGBColor(0, 0, 0))

# -----------------------------
# 36. Total Assets by Asset Group
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Total Assets by Asset Group"

chart_data = CategoryChartData()
chart_data.categories = cap_df['asset_group']
chart_data.add_series('Quantity', cap_df['total_assets'])

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4),
    chart_data
).chart

# 37. Apply colors correctly
for series in chart.series:
    for j, point in enumerate(series.points):
        category_name = chart_data.categories[j].label
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = asset_colors.get(category_name, RGBColor(0, 0, 0))

# -----------------------------
# 38. Annual Depreciation per Asset
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Annual Depreciation per Asset"

chart_data = CategoryChartData()
chart_data.categories = dep_df['asset_name']
chart_data.add_series('Annual Depreciation', dep_df['annual_depreciation'])

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4),
    chart_data
).chart

# 39. Optional: single color for depreciation
for series in chart.series:
    for point in series.points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor(91, 155, 213)  # Blue

# =================================================
# 40. Actuals Comparison
# =================================================
merged = pd.merge(
    actual_2025,
    actual_2026,
    on="account_name",
    how="left"
)

merged = clean_df(merged)

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Actuals Comparison (2025 vs 2026)"

chart_data = CategoryChartData()
chart_data.categories = merged["account_name"].tolist()

chart_data.add_series(
    "2025 Actual",
    merged["total_2025"].tolist()
)

chart_data.add_series(
    "2026 Actual",
    merged["total_2026"].tolist()
)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4),
    chart_data
).chart


# =================================================
# 41. Growth Trend
# =================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Growth Trend Analysis"

chart_data = CategoryChartData()
chart_data.categories = growth_df["account_name"].tolist()

chart_data.add_series(
    "Growth Rate",
    growth_df["growth_rate"].tolist()
)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4),
    chart_data
).chart


# =================================================
# 42. Budget Forecast
# =================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Budget Forecast 2027"

chart_data = CategoryChartData()
chart_data.categories = forecast_df["account_name"].tolist()

chart_data.add_series(
    "Forecast 2027",
    forecast_df["forecast_2027"].tolist()
)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4),
    chart_data
).chart


# =================================================
# 43. Executive Financial Trend
# =================================================

trend_df = pd.merge(
    actual_2025[["account_name","total_2025"]],
    forecast_df[["account_name","forecast_2027"]],
    on="account_name",
    how="left"
)

trend_df = clean_df(trend_df)

revenue_2025 = trend_df.loc[
    trend_df["account_name"]=="Sales Revenue",
    "total_2025"
].sum()

revenue_2027 = trend_df.loc[
    trend_df["account_name"]=="Sales Revenue",
    "forecast_2027"
].sum()

expense_2025 = trend_df.loc[
    trend_df["account_name"]!="Sales Revenue",
    "total_2025"
].sum()

expense_2027 = trend_df.loc[
    trend_df["account_name"]!="Sales Revenue",
    "forecast_2027"
].sum()

profit_2025 = revenue_2025 - expense_2025
profit_2027 = revenue_2027 - expense_2027


slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Revenue vs Expense vs Profit Trend"

chart_data = CategoryChartData()
chart_data.categories = ["2025", "2027 Forecast"]

chart_data.add_series(
    "Revenue",
    [revenue_2025, revenue_2027]
)

chart_data.add_series(
    "Expenses",
    [expense_2025, expense_2027]
)

chart_data.add_series(
    "Profit",
    [profit_2025, profit_2027]
)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1),
    Inches(1.5),
    Inches(8),
    Inches(4),
    chart_data
).chart

# =====================================================
# 44. Monthly Expense Variance (Defined & Labeled)
# =====================================================
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Monthly Expense Variance by Account"

chart_data = CategoryChartData()
# This creates the X-axis labels (e.g., "1 - Salaries")
exp_variance['label'] = exp_variance['month'].astype(str) + " - " + exp_variance['account_name']
chart_data.categories = exp_variance['label'].tolist()

# The first string in add_series defines the Legend label
chart_data.add_series('Actual Expense', [abs(x) for x in exp_variance['total_expense']])
chart_data.add_series('Monthly Variance', [abs(x) for x in exp_variance['variance']])

chart_graphic = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.2),
    Inches(9), Inches(5),
    chart_data
)
chart = chart_graphic.chart

# ENABLE DEFINITIONS
chart.has_legend = True  # Tells the user what Blue and Orange mean
chart.legend.include_in_layout = False

# Series 0: Actual (Blue)
s0 = chart.series[0]
s0.format.fill.solid()
s0.format.fill.fore_color.rgb = RGBColor(91, 155, 213)
s0.has_data_labels = True
s0.data_labels.number_format = '#,##0'

# Series 1: Variance (Orange)
s1 = chart.series[1]
s1.format.fill.solid()
s1.format.fill.fore_color.rgb = RGBColor(237, 125, 49)
s1.has_data_labels = True
s1.data_labels.number_format = '#,##0'

# =====================================================
# 45. Overall Monthly Expense Trend (Positive Line)
# =====================================================
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Overall Monthly Expense Trend"

chart_data = CategoryChartData()
# X-Axis: Months (1, 2, 3...)
chart_data.categories = exp_trend['month'].astype(str).tolist()

# Y-Axis: Monthly Total (Converted to Positive for clear trending)
chart_data.add_series('Monthly Total Expense', [abs(x) for x in exp_trend['monthly_expense']])

chart_graphic = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(1), Inches(1.5),
    Inches(8), Inches(4),
    chart_data
)
chart = chart_graphic.chart

# 1. ENFORCE COLOR AND STYLE
series = chart.series[0]
series.format.line.color.rgb = RGBColor(68, 114, 196) # Professional Blue
series.format.line.width = Inches(0.05)              # Thicker line for visibility

# 2. ADD DATA POINTS & LABELS
series.has_data_labels = True
series.data_labels.number_format = '#,##0'
series.data_labels.font.size = Inches(0.12)

# 3. ADD MARKERS (Dots on the line)
series.marker.style = 2 # Diamond or Circle style
series.marker.size = 7

# ===============================
# 46. Journal Impact (Clustered Column)
# ===============================
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Journal Entry Impact Analysis"

chart_data = CategoryChartData()
chart_data.categories = exp_comparison['month'].astype(str).tolist()

# Add series with numeric conversion
chart_data.add_series('System Entries (Normal)', [float(abs(x)) for x in exp_comparison['expense_without_journal']])
chart_data.add_series('Total Incl. Journals', [float(abs(x)) for x in exp_comparison['expense_with_journal']])
chart_data.add_series('Adjustment Impact', [float(abs(x)) for x in exp_comparison['journal_impact']])

chart_graphic = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.5),
    Inches(9), Inches(4.5),  # width x height
    chart_data
)
chart = chart_graphic.chart
chart.has_legend = True

# Colors and safe data labels (inside for small bars)
colors = [RGBColor(165, 165, 165), RGBColor(91, 155, 213), RGBColor(255, 192, 0)]
for i, series in enumerate(chart.series):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = colors[i]
    series.has_data_labels = True
    series.data_labels.position = XL_LABEL_POSITION.INSIDE_END
    series.data_labels.font.size = Pt(8)

# Rotate X-axis labels
category_axis_46 = chart.category_axis
category_axis_46.tick_labels.rotation = 45

# ===============================
# 47. Monthly Depreciation Forecast (Stacked Column) - FIXED
# ===============================
slide = prs.slides.add_slide(blank_layout)
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "Monthly Depreciation Forecast"

pivot_depr = depr_forecast.pivot(index='forecast_month', columns='asset_group', values='total_depreciation').fillna(0)

chart_data = CategoryChartData()
chart_data.categories = [str(x) for x in pivot_depr.index]

# Add each asset group as a series
for group in pivot_depr.columns:
    chart_data.add_series(group, [float(x) for x in pivot_depr[group]])

chart_graphic = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED,
    Inches(0.5), Inches(1.5),
    Inches(9), Inches(4.5),
    chart_data
)
chart = chart_graphic.chart
chart.has_legend = True

# FIX: Set labels to CENTER or INSIDE_BASE for stacked charts
# 'OUTSIDE_END' is what caused the file corruption/disappearance
for series in chart.series:
    series.has_data_labels = True
    series.data_labels.position = XL_LABEL_POSITION.CENTER # Changed from OUTSIDE_END
    series.data_labels.font.size = Pt(8)
    series.data_labels.number_format = '#,##0'

# Rotate X-axis labels for better fit
chart.category_axis.tick_labels.rotation = 45


# ===============================
# 48. CWIP Pie Chart (with validation)
# ===============================
# Only create pie chart if values exist
if cwip_summary['cwip_value'].sum() > 0:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5)).text = "CWIP Value by Asset Group"

    chart_data = CategoryChartData()
    chart_data.categories = cwip_summary['asset_group'].tolist()
    chart_data.add_series('CWIP Amount', [float(x) for x in cwip_summary['cwip_value']])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(2), Inches(1.5),
        Inches(6), Inches(4.5),
        chart_data
    ).chart
    chart.has_legend = True

    # Show percentage labels outside pie slices
    chart.series[0].has_data_labels = True
    chart.series[0].data_labels.show_percentage = True
    chart.series[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.series[0].data_labels.font.size = Pt(10)

# -----------------------------
# Save file
# -----------------------------
prs.save("Auto_Financial_Pack_2025.pptx")

print("Financial Pack Generated Successfully")